"""Build the supervisor deck as a .pptx for Google Slides.

Text is deliberately sparse: a slide is a prompt for the speaker, not a
document. Every number here is the frozen value from reference/*.json.

Fonts are restricted to Georgia (headlines) and Arial (everything else)
because Google Slides has both natively -- anything else silently
substitutes on upload and wrecks the layout.
"""
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

C = {
    "ink":   RGBColor(0x17, 0x1A, 0x21),
    "soft":  RGBColor(0x4A, 0x52, 0x60),
    "faint": RGBColor(0x79, 0x82, 0x8F),
    "rule":  RGBColor(0xDD, 0xE2, 0xEA),
    "accent":RGBColor(0x2F, 0x5D, 0x8C),
    "good":  RGBColor(0x1B, 0x7A, 0x4B),
    "bad":   RGBColor(0xB2, 0x3A, 0x34),
    "warn":  RGBColor(0x9A, 0x6A, 0x12),
    "white": RGBColor(0xFF, 0xFF, 0xFF),
    "panel": RGBColor(0xF5, 0xF7, 0xFA),
}
SERIF, SANS = "Georgia", "Arial"
W, H = Inches(13.333), Inches(7.5)
M = Inches(0.62)                      # page margin
CASES = "reference/cases"

prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]


def tb(slide, x, y, w, h, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    tf.paragraphs[0].alignment = align
    return tf


def run(p, text, size, color, font=SANS, bold=False, italic=False, spacing=None):
    r = p.add_run(); r.text = text
    r.font.size, r.font.name, r.font.bold, r.font.italic = Pt(size), font, bold, italic
    r.font.color.rgb = color
    return r


def eyebrow(slide, text, color):
    tf = tb(slide, M, Inches(0.42), Inches(9), Inches(0.3))
    run(tf.paragraphs[0], text.upper(), 11, color, bold=True)


def title(slide, text, y=Inches(0.78), size=31, width=11.4):
    tf = tb(slide, M, y, Inches(width), Inches(1.0))
    run(tf.paragraphs[0], text, size, C["ink"], font=SERIF, bold=True)


def slide_num(slide, n, total=8):
    tf = tb(slide, W - M - Inches(1.2), H - Inches(0.52), Inches(1.2), Inches(0.3),
            align=PP_ALIGN.RIGHT)
    run(tf.paragraphs[0], f"{n} / {total}", 10, C["faint"])


def context(slide, dataset, arm=None):
    """Bottom-left label naming the dataset and which task the slide is about.

    Added after a read-through showed the deck never said there were two
    tasks at all -- so perception's 0.835 and reasoning's 0.520 looked like
    contradictory numbers for the same thing, and one slide silently mixed
    both arms.
    """
    tf = tb(slide, M, H - Inches(0.54), Inches(9), Inches(0.32))
    p = tf.paragraphs[0]
    run(p, dataset, 10, C["accent"], bold=True)
    if arm:
        run(p, "   ·   ", 10, C["rule"])
        run(p, arm, 10, C["faint"])


def rule(slide, y, x=M, w=None):
    w = w or (W - 2 * M)
    ln = slide.shapes.add_shape(1, x, y, w, Emu(9525))   # 1 = rectangle
    ln.fill.solid(); ln.fill.fore_color.rgb = C["rule"]
    ln.line.fill.background(); ln.shadow.inherit = False
    return ln


def panel(slide, x, y, w, h, fill=None, accent=None):
    sh = slide.shapes.add_shape(1, x, y, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill or C["panel"]
    sh.line.color.rgb = C["rule"]; sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    if accent:
        bar = slide.shapes.add_shape(1, x, y, w, Inches(0.045))
        bar.fill.solid(); bar.fill.fore_color.rgb = accent
        bar.line.fill.background(); bar.shadow.inherit = False
    return sh


def table(slide, x, y, w, rows, col_w, header=None, row_h=Inches(0.42)):
    """Minimal table: header row + body, hairline separators, no grid."""
    yy = y
    if header:
        cx = x
        for text, cw in zip(header, col_w):
            tf = tb(slide, cx, yy, cw, Inches(0.26))
            run(tf.paragraphs[0], text.upper(), 9.5, C["faint"], bold=True)
            cx += cw
        yy += Inches(0.30)
        rule(slide, yy); yy += Inches(0.12)
    for r in rows:
        cx = x
        for cell, cw in zip(r, col_w):
            text, *style = cell if isinstance(cell, tuple) else (cell,)
            color = style[0] if style else C["ink"]
            bold = style[1] if len(style) > 1 else False
            size = style[2] if len(style) > 2 else 13
            ci = style[3] if len(style) > 3 else None
            tf = tb(slide, cx, yy, cw, row_h, anchor=MSO_ANCHOR.TOP)
            run(tf.paragraphs[0], text, size, color, bold=bold)
            if ci:
                # CI rides in the same cell, smaller and grey, so the point
                # estimate stays scannable while the interval is always
                # visible next to it -- never on its own line, which reads
                # as a separate number.
                run(tf.paragraphs[0], "  " + ci, size - 3.5, C["faint"])
            cx += cw
        yy += row_h
        rule(slide, yy - Inches(0.08))
    return yy


# ───────────────────── 1 · FERMAT + the two tasks ─────────────────────
s = prs.slides.add_slide(BLANK)
eyebrow(s, "The dataset", C["accent"])
title(s, "FERMAT: handwritten math,\nerrors planted on purpose")

s.shapes.add_picture(f"{CASES}/07_qwen7b_confidently_wrong_grading/image.jpg",
                     M, Inches(2.3), height=Inches(3.35))

x2 = Inches(6.2)
tf = tb(s, x2, Inches(2.25), Inches(6.5), Inches(1.0))
for i, (a, b) in enumerate([("~2,200", " photographed solutions   ·   "),
                            ("85%", " contain a deliberate error")]):
    p = tf.paragraphs[0] if i == 0 else p
    run(p, a, 15, C["ink"], bold=True); run(p, b, 15, C["soft"])

# The two arms, stated up front -- everything downstream is one or the other.
panel(s, x2, Inches(2.85), Inches(6.5), Inches(1.85), accent=C["accent"])
tf = tb(s, x2 + Inches(0.28), Inches(3.08), Inches(5.95), Inches(1.5))
run(tf.paragraphs[0], "WE TEST TWO SEPARATE TASKS", 10, C["accent"], bold=True)
p = tf.add_paragraph(); p.space_before = Pt(9)
run(p, "Perception", 14, C["ink"], bold=True)
run(p, " — read the page. Did it transcribe correctly?", 14, C["soft"])
p = tf.add_paragraph(); p.space_before = Pt(6)
run(p, "Reasoning", 14, C["ink"], bold=True)
run(p, " — grade the page. Did it spot the error correctly?", 14, C["soft"])

tf = tb(s, x2, Inches(5.0), Inches(6.5), Inches(0.3))
run(tf.paragraphs[0], "ONLY FERMAT HAS CORRECT ANSWERS TOO", 10, C["faint"], bold=True)
table(s, x2, Inches(5.38), Inches(6.5),
      rows=[["FERMAT", ("15% clean", C["good"], True)],
            ["ErrorRadar  /  ScratchMath", ("no clean items", C["bad"], True)]],
      col_w=[Inches(3.6), Inches(2.9)], row_h=Inches(0.4))
context(s, "FERMAT", "setup")
slide_num(s, 1, 8)

# ───────────────────── 2 · perception result ─────────────────────
s = prs.slides.add_slide(BLANK)
eyebrow(s, "Perception — confirmed", C["good"])
title(s, "When it can't read a page twice the same way, it's wrong")

tf = tb(s, M, Inches(1.9), Inches(11.9), Inches(0.85))
p0 = tf.paragraphs[0]
run(p0, "The model is asked to ", 15, C["soft"])
run(p0, "\u201cexplicitly perform OCR on the handwritten text and extract the "
        "content in LaTeX format.\u201d", 15, C["ink"], italic=True)
p0 = tf.add_paragraph(); p0.space_before = Pt(5)
run(p0, "We do that 5 times per page and measure how much the readings disagree.",
    15, C["soft"])

table(s, M, Inches(2.72), Inches(6.9),
      header=["", "AUROC  (95% CI)"],
      rows=[[("All 300 pages", C["ink"], True),
             ("0.835", C["good"], True, 15, "[0.787, 0.879]")],
            [("Hardest cases deleted  (slide 3)", C["ink"], True),
             ("0.796", C["good"], True, 15, "[0.736, 0.852]")],
            [("The model's own confidence", C["ink"], True),
             ("0.537", C["bad"], True, 15, "[0.469, 0.605]")]],
      col_w=[Inches(3.9), Inches(3.0)], row_h=Inches(0.5))

tf = tb(s, M, Inches(4.5), Inches(6.5), Inches(0.6))
run(tf.paragraphs[0], "Row 3's interval crosses 0.50 — the model's own confidence "
    "score is no better than a coin flip here.", 12, C["soft"], italic=True)

panel(s, Inches(7.9), Inches(2.6), Inches(4.8), Inches(2.2), accent=C["good"])
tf = tb(s, Inches(8.22), Inches(2.88), Inches(4.2), Inches(1.9))
run(tf.paragraphs[0], "All 5 readings disagree?", 18, C["ink"], font=SERIF, bold=True)
p = tf.add_paragraph(); p.space_before = Pt(8)
run(p, "Wrong 57 times out of 61.", 18, C["good"], font=SERIF, bold=True)
p = tf.add_paragraph(); p.space_before = Pt(12)
run(p, "Flag those for a human. Catches 57 real errors, wastes 4.", 12.5, C["soft"])
# Cross-model strip: the perception story including the trap, in three
# rows. LLaVA is folded into the caption rather than given a row of
# dashes -- it never produced a measurement, so it is context, not a
# comparison. Added after Pixtral replicated the result on 2026-08-09.
panel(s, M, Inches(5.32), W - 2 * M, Inches(1.5), accent=C["accent"])
tf = tb(s, M + Inches(0.3), Inches(5.5), Inches(3.6), Inches(1.25))
run(tf.paragraphs[0], "ACROSS MODEL FAMILIES", 10, C["accent"], bold=True)
p3 = tf.add_paragraph(); p3.space_before = Pt(5)
run(p3, "A high score means nothing until the control is applied.",
    12.5, C["soft"])
p3 = tf.add_paragraph(); p3.space_before = Pt(4)
run(p3, "A fourth, LLaVA-NeXT-7B, could not read the pages at all.",
    11.5, C["faint"])

table(s, Inches(4.8), Inches(5.5), Inches(7.9),
      header=["", "AUROC", "after control", ""],
      rows=[[("Qwen2.5-VL-3B", C["ink"], True), ("0.835", C["good"], True, 13),
             ("0.796", C["good"], True, 13), ("holds", C["good"], False, 12)],
            [("Pixtral-12B", C["ink"], True), ("0.828", C["good"], True, 13),
             ("0.772", C["good"], True, 13), ("holds", C["good"], False, 12)],
            [("InternVL3-8B", C["ink"], True), ("0.915", C["bad"], True, 13),
             ("0.556", C["bad"], True, 13), ("collapses", C["bad"], False, 12)]],
      col_w=[Inches(2.4), Inches(1.4), Inches(1.9), Inches(2.2)], row_h=Inches(0.29))

context(s, "FERMAT", "perception — reading the page")
slide_num(s, 2, 8)

# ─────────────── 3 · perception evidence + the check ───────────────
s = prs.slides.add_slide(BLANK)
eyebrow(s, "Perception — evidence", C["accent"])
title(s, "Where it works, and the blind spot", size=30)

quad = [
    ("01_qwen3b_perception_low_entropy_correct", "Agrees  ·  H = 0.00", "Correct", C["good"]),
    ("02_qwen3b_perception_high_entropy_wrong",  "Disagrees  ·  H = 1.61", "Caught", C["good"]),
    ("03_qwen3b_perception_high_entropy_correct","Disagrees  ·  H = 1.61", "False alarm", C["warn"]),
    ("04_qwen3b_perception_low_entropy_wrong",   "Agrees  ·  H = 0.00", "Blind spot", C["bad"]),
]
cw, ch = Inches(3.02), Inches(2.95)
BOX_W, BOX_H = cw - Inches(0.3), Inches(1.95)
for i, (case, lab, verdict, col) in enumerate(quad):
    cx = M + (cw + Inches(0.14)) * i
    panel(s, cx, Inches(1.9), cw, ch, fill=C["white"], accent=col)
    with Image.open(f"{CASES}/{case}/image.jpg") as im:
        iw, ih = im.size
    sc = min(BOX_W / iw, BOX_H / ih)
    pw, ph = int(iw * sc), int(ih * sc)
    s.shapes.add_picture(f"{CASES}/{case}/image.jpg",
                         cx + (cw - pw) // 2, Inches(2.06) + (BOX_H - ph) // 2,
                         width=pw, height=ph)
    tf = tb(s, cx + Inches(0.15), Inches(4.16), cw - Inches(0.3), Inches(0.6))
    run(tf.paragraphs[0], lab, 10.5, C["soft"])
    p = tf.add_paragraph(); p.space_before = Pt(2)
    run(p, verdict, 13, col, bold=True)

panel(s, M, Inches(5.05), W - 2 * M, Inches(1.62), accent=C["bad"])
tf = tb(s, M + Inches(0.32), Inches(5.28), Inches(11.85), Inches(1.35))
run(tf.paragraphs[0], "Page 1 and page 4 look identical to the method — both entropy zero. "
    "One is right, one is wrong.", 17, C["ink"], font=SERIF, bold=True)
p = tf.add_paragraph(); p.space_before = Pt(9)
run(p, "The robustness check: ", 13.5, C["accent"], bold=True)
run(p, "delete every page where all 5 readings differed — easy to flag, nearly "
       "always wrong. If the score survives, the signal is graded. Qwen ",
    13.5, C["soft"])
run(p, "0.835 \u2192 0.794", 13.5, C["good"], bold=True)
run(p, ", 239 pages left.", 13.5, C["soft"])
context(s, "FERMAT", "perception — Qwen2.5-VL-3B")
slide_num(s, 3, 8)

# --------------- 4 . reasoning does not work ---------------
s = prs.slides.add_slide(BLANK)
eyebrow(s, "Reasoning \u2014 negative", C["bad"])
title(s, "The same signal does not work for grading")

tf = tb(s, M, Inches(1.8), Inches(11.7), Inches(0.85))
p0 = tf.paragraphs[0]
run(p0, "Different task. The model is asked only: ", 15, C["soft"])
run(p0, "\u201canalyze the Answer to determine whether there is any error.\u201d",
    15, C["ink"], italic=True)
run(p0, "  It never sees the correct answer.", 15, C["soft"])

table(s, M, Inches(2.95), Inches(11),
      header=["", "AUROC  (95% CI)", ""],
      rows=[[("Grading, 300 balanced pages", C["ink"], True),
             ("0.520", C["bad"], True, 17, "[0.458, 0.582]"),
             ("interval contains chance", C["bad"], False, 13)]],
      col_w=[Inches(4.6), Inches(3.4), Inches(3.0)], row_h=Inches(0.7))

panel(s, M, Inches(4.35), W - 2 * M, Inches(1.95), accent=C["bad"])
tf = tb(s, M + Inches(0.35), Inches(4.62), Inches(11.8), Inches(1.6))
run(tf.paragraphs[0], "No signal. That is the honest reasoning result.",
    22, C["ink"], font=SERIF, bold=True)
p2 = tf.add_paragraph(); p2.space_before = Pt(10)
run(p2, "The model answers \u201cerror\u201d for 80\u201390% of pages regardless of "
        "what is on them. Balancing the set is what makes this readable \u2014 on an "
        "unbalanced set that bias alone would score well.", 14, C["soft"])
context(s, "FERMAT", "reasoning \u2014 grading the page  \u00b7  Qwen2.5-VL-7B")
slide_num(s, 4, 8)

# --------------- 5 . the trap ---------------
s = prs.slides.add_slide(BLANK)
eyebrow(s, "What we got wrong", C["warn"])
title(s, "We believed the opposite for three weeks", size=30)

tf = tb(s, M, Inches(1.72), Inches(11.9), Inches(0.6))
run(tf.paragraphs[0], "Splitting the grading data by ground truth seemed to "
    "recover a strong effect \u2014 replicated across three model families.",
    15, C["soft"])

table(s, M, Inches(2.45), Inches(11.6),
      header=["Model", "We reported", "Signal-free coin", "Collapse"],
      rows=[[("Qwen2.5-VL-3B", C["ink"], True), ("0.854", C["good"], True, 15),
             ("0.901", C["bad"], True, 15), ("99.8%", C["bad"], True, 14)],
            [("Qwen2.5-VL-7B", C["ink"], True), ("0.801", C["good"], True, 15),
             ("0.868", C["bad"], True, 15), ("100%", C["bad"], True, 14)],
            [("LLaVA-NeXT-7B", C["ink"], True), ("0.775", C["good"], True, 15),
             ("0.831", C["bad"], True, 15), ("97.2%", C["bad"], True, 14)]],
      col_w=[Inches(3.4), Inches(2.7), Inches(3.0), Inches(2.5)], row_h=Inches(0.5))

panel(s, M, Inches(4.7), W - 2 * M, Inches(1.9), accent=C["bad"])
tf = tb(s, M + Inches(0.35), Inches(4.95), Inches(11.8), Inches(1.6))
run(tf.paragraphs[0], "Inside one label group, \u201cwas it correct\u201d IS "
    "\u201cwhat did it answer\u201d.", 19, C["ink"], font=SERIF, bold=True)
p2 = tf.add_paragraph(); p2.space_before = Pt(8)
run(p2, "The same column, 100% of the time. Entropy is computed from those very "
        "votes, so it predicts correctness almost by definition. ", 13.5, C["soft"])
run(p2, "A coin with no signal scores higher than we did.", 13.5, C["bad"], bold=True)
p2 = tf.add_paragraph(); p2.space_before = Pt(6)
run(p2, "Not caught by pre-registration, by power, or by replicating it three times.",
    13, C["faint"], italic=True)
context(s, "FERMAT", "reasoning \u2014 retracted 2026-08-09")
slide_num(s, 5, 8)

# ─────────────── 6 · what did not work ───────────────
s = prs.slides.add_slide(BLANK)
eyebrow(s, "What did not work", C["bad"])
title(s, "Two results we had to throw out")

# InternVL3 -- perception, failed the robustness check
panel(s, M, Inches(2.05), Inches(5.95), Inches(3.5), accent=C["bad"])
tf = tb(s, M + Inches(0.32), Inches(2.3), Inches(5.3), Inches(3.1))
run(tf.paragraphs[0], "A 4TH MODEL, PERCEPTION", 10, C["faint"], bold=True)
p = tf.add_paragraph(); p.space_before = Pt(9)
run(p, "0.915", 27, C["bad"], font=SERIF, bold=True)
run(p, "  \u2192  ", 18, C["faint"])
run(p, "0.556", 27, C["ink"], font=SERIF, bold=True)
p = tf.add_paragraph(); p.space_before = Pt(11)
run(p, "InternVL3 looked like our best result. Under the robustness check it "
       "collapsed to chance — ", 13.5, C["soft"])
run(p, "2 of every 3 pages were total breakdowns", 13.5, C["ink"], bold=True)
run(p, ", and that binary was all the score was detecting.", 13.5, C["soft"])

# ScratchMath -- the second dataset
panel(s, Inches(6.95), Inches(2.05), Inches(5.75), Inches(3.5), accent=C["warn"])
tf = tb(s, Inches(7.27), Inches(2.3), Inches(5.15), Inches(3.1))
run(tf.paragraphs[0], "A 2ND DATASET, REASONING", 10, C["faint"], bold=True)
p = tf.add_paragraph(); p.space_before = Pt(9)
run(p, "96%", 27, C["warn"], font=SERIF, bold=True)
run(p, "  accuracy, and meaningless", 14, C["soft"])
p = tf.add_paragraph(); p.space_before = Pt(11)
run(p, "ScratchMath, 100 pages. Every page contains an error, and the model "
       "answers “error” 90% of the time — so it scores well by construction.",
    13.5, C["soft"])
p = tf.add_paragraph(); p.space_before = Pt(8)
run(p, "In 24% of readings it said it could not read the image. "
       "It answered “error” anyway in 82% of those.", 13.5, C["ink"], bold=True)

panel(s, M, Inches(5.78), W - 2 * M, Inches(0.95), accent=C["accent"])
tf = tb(s, M + Inches(0.35), Inches(5.98), Inches(11.8), Inches(0.6))
run(tf.paragraphs[0], "Both were caught by checks we run on every number — "
    "which is why the results on the earlier slides can be trusted.",
    14.5, C["ink"], font=SERIF, bold=True)
context(s, "InternVL3 on FERMAT  ·  Qwen2.5-VL-7B on ScratchMath")
slide_num(s, 6, 8)

# ─────────────── 7 · open decision ───────────────
s = prs.slides.add_slide(BLANK)
eyebrow(s, "Open decision", C["accent"])
title(s, "One question")

panel(s, M, Inches(2.0), Inches(5.95), Inches(2.75), accent=C["good"])
tf = tb(s, M + Inches(0.32), Inches(2.3), Inches(5.3), Inches(2.3))
run(tf.paragraphs[0], "SOLID", 11, C["good"], bold=True)
for t in ["Perception: 0.835, replicated at 0.828\n    on a second model family",
          "Abstention rule: 57 of 61",
          "Reasoning: no signal \u2014 measured, not assumed",
          "The stratification trap: a method finding"]:
    p = tf.add_paragraph(); p.space_before = Pt(9)
    run(p, "— " + t, 13.5, C["ink"])

panel(s, Inches(6.95), Inches(2.0), Inches(5.75), Inches(2.75), accent=C["warn"])
tf = tb(s, Inches(7.27), Inches(2.3), Inches(5.1), Inches(2.3))
run(tf.paragraphs[0], "LIMITED", 11, C["warn"], bold=True)
for t in ["Every result sits on FERMAT alone",
          "Both alternative datasets are 100%\n    error items",
          "…so neither can test the split on\n    slide 4 at all"]:
    p = tf.add_paragraph(); p.space_before = Pt(9)
    run(p, "— " + t, 13.5, C["ink"])

panel(s, M, Inches(5.05), W - 2 * M, Inches(1.5), accent=C["accent"])
tf = tb(s, M + Inches(0.35), Inches(5.35), Inches(11.8), Inches(1.0))
run(tf.paragraphs[0], "Accept one dataset as a limitation — or label correct pages ourselves?",
    21, C["ink"], font=SERIF, bold=True)
p = tf.add_paragraph(); p.space_before = Pt(8)
run(p, "~3 weeks to submission. All experiments done; remaining work is the write-up.",
    13.5, C["soft"])
context(s, "FERMAT", "both tasks")
slide_num(s, 7, 8)

# ─────────────── 8 · the experiment we would still run ───────────────
s = prs.slides.add_slide(BLANK)
eyebrow(s, "Proposed next experiment", C["accent"])
title(s, "Make it solve the problem, not just eyeball it")

panel(s, M, Inches(1.95), W - 2 * M, Inches(1.25), accent=C["warn"])
tf = tb(s, M + Inches(0.32), Inches(2.18), Inches(11.85), Inches(1.0))
run(tf.paragraphs[0], "The gap: plausibility-checking instead of recomputing.",
    18, C["ink"], font=SERIF, bold=True)
p = tf.add_paragraph(); p.space_before = Pt(7)
run(p, "It reads the student's steps and judges whether they look right. "
       "One page asserts ", 13.5, C["soft"])
run(p, "19 \u00d7 \u221218 \u00d7 23 = \u22127966", 13.5, C["ink"], bold=True)
run(p, ". The model accepted it. The correct value is ", 13.5, C["soft"])
run(p, "\u22127866", 13.5, C["bad"], bold=True)
run(p, ".", 13.5, C["soft"])

# what has already been screened
panel(s, M, Inches(3.45), Inches(5.95), Inches(2.7), accent=C["bad"])
tf = tb(s, M + Inches(0.32), Inches(3.68), Inches(5.3), Inches(2.2))
run(tf.paragraphs[0], "ALREADY TRIED — PROMPT LEVEL", 10, C["faint"], bold=True)
p = tf.add_paragraph(); p.space_before = Pt(8)
run(p, "\u201cFirst, transcribe the final result the student arrived at, exactly "
       "as written. Then check that result by working the problem yourself. "
       "Only then decide whether the Answer contains an error.\u201d",
    12, C["ink"], italic=True)
p = tf.add_paragraph(); p.space_before = Pt(9)
run(p, "Qwen-3B  ", 13, C["ink"], bold=True)
run(p, "says-error 94% \u2192 38%, accuracy 0.51 \u2192 0.55", 13, C["soft"])
p = tf.add_paragraph(); p.space_before = Pt(5)
run(p, "InternVL3  ", 13, C["ink"], bold=True)
run(p, "accuracy 0.70 \u2192 0.67", 13, C["soft"])
p = tf.add_paragraph(); p.space_before = Pt(9)
run(p, "Behaviour swung hard. Accuracy did not move.", 13, C["bad"], bold=True)

# what is still open
panel(s, Inches(6.95), Inches(3.45), Inches(5.75), Inches(2.7), accent=C["good"])
tf = tb(s, Inches(7.27), Inches(3.68), Inches(5.15), Inches(2.2))
run(tf.paragraphs[0], "NOT YET TRIED — PIPELINE LEVEL", 10, C["faint"], bold=True)
p = tf.add_paragraph(); p.space_before = Pt(8)
run(p, "Solve the page independently, then compare in code.", 13.5, C["ink"], bold=True)
p = tf.add_paragraph(); p.space_before = Pt(9)
run(p, "Asking one prompt to do both still lets the model see the student's "
       "answer first, so it can anchor on it. A separate solve step cannot.",
    13, C["soft"])
p = tf.add_paragraph(); p.space_before = Pt(8)
run(p, "Then re-measure: does the split on slide 4 still hold, and does "
       "uncertainty still predict the errors?", 13, C["soft"])

panel(s, M, Inches(6.22), W - 2 * M, Inches(0.6), accent=C["accent"])
tf = tb(s, M + Inches(0.35), Inches(6.36), Inches(11.8), Inches(0.42))
run(tf.paragraphs[0], "Worth the remaining time — or write up what we have?",
    15, C["ink"], font=SERIF, bold=True)
context(s, "FERMAT", "reasoning — proposed")
slide_num(s, 8, 8)

out = "slides/fermat_findings.pptx"
prs.save(out)
import os
print(f"saved {out}  ({os.path.getsize(out)/1024/1024:.2f} MB, {len(prs.slides._sldIdLst)} slides)")
